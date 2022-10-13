#!/usr/bin/python3

from pptx import Presentation
from glob import glob
import random
import sys
from re import sub


SLIDES_COUNT = 10


def snake_case(s):
    s = sub(r'[^a-zA-Z0-9]+', '_', s)
    return '_'.join(
        sub('([A-Z][a-z]+)', r' \1',
        sub('([A-Z]+)', r' \1',
        s.replace('-', ' '))).split()).lower()


presentation_title = sys.argv[1]
print(f"Generating presentation for {presentation_title}")

presentation = Presentation('template_presentation.pptx')

title_layout = presentation.slide_layouts[0]
blank_layout = presentation.slide_layouts[6]

first_slide = presentation.slides.add_slide(title_layout)


first_slide.shapes.title.text = presentation_title
first_slide.placeholders[1].text = "Powerpoint Roulette"

# Select a random images
all_image_filenames = glob('images/*')
image_filenames = random.sample(all_image_filenames, k=SLIDES_COUNT)

for image_filename in image_filenames:
    image_slide = presentation.slides.add_slide(blank_layout)

    # Will strecth or shrink the image to the slide size
    pic = image_slide.shapes.add_picture(
        image_filename, 0, 0, 
        presentation.slide_width, presentation.slide_height)

last_slide = presentation.slides.add_slide(title_layout)
last_slide.shapes.title.text = presentation_title
last_slide.placeholders[1].text = "End. Well done!"

filename = snake_case(presentation_title) + "_presentation.pptx"
presentation.save(filename)
